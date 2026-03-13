from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable, Sequence

import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = REPO_ROOT / "docs"
FIG_DIR = DOCS_DIR / "figures" / "progress_report_20260312"
PPT_PATH = DOCS_DIR / "progress_report_20260312_en.pptx"


PALETTE = {
    "navy": (25, 47, 89),
    "blue": (58, 110, 165),
    "sky": (102, 163, 224),
    "teal": (48, 128, 136),
    "green": (70, 143, 83),
    "orange": (201, 117, 55),
    "red": (188, 74, 72),
    "ink": (34, 43, 53),
    "muted": (108, 118, 128),
    "paper": (247, 248, 250),
    "white": (255, 255, 255),
}


def _load_font(size: int, *, bold: bool = False) -> ImageFont.ImageFont:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size=size)
    return ImageFont.load_default()


FONT_H1 = _load_font(42, bold=True)
FONT_H2 = _load_font(30, bold=True)
FONT_BODY = _load_font(22)
FONT_SMALL = _load_font(18)


def _wrap(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = word if not current else f"{current} {word}"
        if draw.textlength(candidate, font=font) <= width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or [text]


def _draw_wrapped_text(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    *,
    font: ImageFont.ImageFont,
    fill: tuple[int, int, int],
    width: int,
    line_gap: int = 8,
) -> int:
    x, y = xy
    for line in _wrap(draw, text, font, width):
        draw.text((x, y), line, font=font, fill=fill)
        y += font.size + line_gap
    return y


def _box(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int, int, int],
    *,
    fill: tuple[int, int, int],
    outline: tuple[int, int, int] = PALETTE["navy"],
    radius: int = 16,
    width: int = 3,
) -> None:
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def _arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], *, fill: tuple[int, int, int], width: int = 5) -> None:
    draw.line([start, end], fill=fill, width=width)
    dx = end[0] - start[0]
    dy = end[1] - start[1]
    if dx == 0 and dy == 0:
        return
    scale = (dx ** 2 + dy ** 2) ** 0.5
    ux, uy = dx / scale, dy / scale
    px, py = -uy, ux
    tip = end
    left = (int(end[0] - 18 * ux + 8 * px), int(end[1] - 18 * uy + 8 * py))
    right = (int(end[0] - 18 * ux - 8 * px), int(end[1] - 18 * uy - 8 * py))
    draw.polygon([tip, left, right], fill=fill)


def make_architecture_figure(path: Path) -> None:
    img = Image.new("RGB", (1600, 900), PALETTE["paper"])
    draw = ImageDraw.Draw(img)
    draw.text((70, 40), "SourceAgent Architecture: Phase A, Phase A.5, and Phase B", font=FONT_H1, fill=PALETTE["navy"])
    draw.text((70, 100), "Deterministic structure recovery remains authoritative. LLMs only assist under bounded contracts.", font=FONT_BODY, fill=PALETTE["muted"])

    boxes = {
        "input": (90, 180, 390, 300),
        "stage12": (450, 165, 820, 320),
        "stage37": (450, 350, 820, 560),
        "stage89": (880, 165, 1235, 380),
        "phasea5": (880, 430, 1235, 640),
        "phaseb": (1290, 235, 1525, 555),
    }

    _box(draw, boxes["input"], fill=(233, 239, 248))
    _box(draw, boxes["stage12"], fill=(221, 235, 247))
    _box(draw, boxes["stage37"], fill=(220, 239, 231))
    _box(draw, boxes["stage89"], fill=(228, 238, 252))
    _box(draw, boxes["phasea5"], fill=(255, 239, 223), outline=PALETTE["orange"])
    _box(draw, boxes["phaseb"], fill=(246, 233, 236), outline=PALETTE["red"])

    draw.text((120, 205), "Input Binary", font=FONT_H2, fill=PALETTE["navy"])
    _draw_wrapped_text(draw, (120, 248), ".elf / stripped ELF / raw .bin", font=FONT_BODY, fill=PALETTE["ink"], width=220)

    draw.text((485, 190), "Stage 1-2", font=FONT_H2, fill=PALETTE["navy"])
    _draw_wrapped_text(
        draw,
        (485, 240),
        "MemoryMap: regions, base, vector table, ISR hypotheses\nMemoryAccessIndex: load/store index, symbol table, decompiled cache",
        font=FONT_BODY,
        fill=PALETTE["ink"],
        width=300,
    )

    draw.text((485, 375), "Stage 3-7", font=FONT_H2, fill=PALETTE["green"])
    _draw_wrapped_text(
        draw,
        (485, 425),
        "Source miners + sink miners\nEvidence packs\nProposal (heuristic or LLM)\nVerifier with label-specific obligations",
        font=FONT_BODY,
        fill=PALETTE["ink"],
        width=300,
    )

    draw.text((915, 190), "Stage 8-9", font=FONT_H2, fill=PALETTE["navy"])
    _draw_wrapped_text(
        draw,
        (915, 240),
        "ChannelGraph + refined objects\nSink-root extraction\nDerive/check summarization\nTunnel-aware chain linking",
        font=FONT_BODY,
        fill=PALETTE["ink"],
        width=280,
    )

    draw.text((910, 455), "Phase A.5", font=FONT_H2, fill=PALETTE["orange"])
    _draw_wrapped_text(
        draw,
        (910, 505),
        "Bounded LLM supervision over low-confidence sources, sinks, objects, and channels.\nAccepted suggestions pass deterministic merge gates and enrich Phase A artifacts.",
        font=FONT_BODY,
        fill=PALETTE["ink"],
        width=280,
    )

    draw.text((1325, 260), "Phase B", font=FONT_H2, fill=PALETTE["red"])
    _draw_wrapped_text(
        draw,
        (1325, 310),
        "Semantic review and verdict calibration.\nDoes not rewrite source reachability, object binding, channel traversal, or root extraction.",
        font=FONT_BODY,
        fill=PALETTE["ink"],
        width=165,
    )

    _arrow(draw, (390, 240), (450, 240), fill=PALETTE["navy"])
    _arrow(draw, (635, 320), (635, 350), fill=PALETTE["blue"])
    _arrow(draw, (820, 260), (880, 260), fill=PALETTE["navy"])
    _arrow(draw, (1235, 270), (1290, 350), fill=PALETTE["red"])
    _arrow(draw, (1058, 430), (1058, 380), fill=PALETTE["orange"])
    _arrow(draw, (1235, 535), (1290, 430), fill=PALETTE["orange"])

    draw.text((1030, 675), "bounded enrichment loop", font=FONT_SMALL, fill=PALETTE["orange"])
    img.save(path)


def make_chain_assembly_figure(path: Path) -> None:
    img = Image.new("RGB", (1600, 900), PALETTE["paper"])
    draw = ImageDraw.Draw(img)
    draw.text((70, 40), "How a Chain Is Assembled", font=FONT_H1, fill=PALETTE["navy"])
    draw.text((70, 100), "The linker does not ask an LLM to guess a chain from scratch. It composes bounded local facts.", font=FONT_BODY, fill=PALETTE["muted"])

    labels = [
        ("1. Source", "attacker-controlled entry\nMMIO / ISR / DMA / shared buffer", (80, 260, 300, 450), (221, 235, 247)),
        ("2. Object", "shared SRAM object or symbol-backed object", (350, 260, 570, 450), (220, 239, 231)),
        ("3. Channel", "cross-context producer → consumer edge if required", (620, 260, 840, 450), (228, 238, 252)),
        ("4. Active Root", "the dangerous value at the sink\nlength / index / format / dispatch", (890, 260, 1110, 450), (255, 239, 223)),
        ("5. Derive / Check", "how the root is computed\nwhether visible checks bind it", (1160, 260, 1380, 450), (246, 233, 236)),
    ]

    for title, body, xy, fill in labels:
        _box(draw, xy, fill=fill, outline=PALETTE["navy"] if fill != (255, 239, 223) else PALETTE["orange"])
        draw.text((xy[0] + 18, xy[1] + 16), title, font=FONT_H2, fill=PALETTE["navy"])
        _draw_wrapped_text(draw, (xy[0] + 18, xy[1] + 78), body, font=FONT_BODY, fill=PALETTE["ink"], width=xy[2] - xy[0] - 36)

    for left, right in zip(labels, labels[1:]):
        _arrow(draw, (left[2][2], 355), (right[2][0], 355), fill=PALETTE["blue"])

    _box(draw, (470, 560, 1130, 760), fill=(236, 244, 250), outline=PALETTE["teal"], radius=22)
    draw.text((505, 590), "A non-drop chain requires more than graph reachability", font=FONT_H2, fill=PALETTE["teal"])
    bullets = [
        "source_reached: a recovered source can still explain the sink-side value",
        "object_bound: the root or data path lands on a recovered object",
        "channel_satisfied: cross-context chains must use a valid channel edge",
        "root_matched: the active sink parameter is the true risk-driving root",
        "derive/check_explained: root derivation and visible guards are explainable",
    ]
    y = 645
    for bullet in bullets:
        draw.ellipse((520, y + 10, 534, y + 24), fill=PALETTE["teal"])
        y = _draw_wrapped_text(draw, (550, y), bullet, font=FONT_BODY, fill=PALETTE["ink"], width=530, line_gap=6) + 10

    draw.text((1260, 605), "Verdict classes", font=FONT_H2, fill=PALETTE["red"])
    verdicts = ["CONFIRMED", "SUSPICIOUS", "SAFE_OR_LOW_RISK", "DROP"]
    colors = [PALETTE["red"], PALETTE["orange"], PALETTE["green"], PALETTE["muted"]]
    y = 655
    for verdict, color in zip(verdicts, colors):
        draw.rectangle((1265, y, 1285, y + 20), fill=color)
        draw.text((1298, y - 4), verdict, font=FONT_BODY, fill=PALETTE["ink"])
        y += 40
    img.save(path)


def make_gt_tiers_figure(path: Path) -> None:
    plt.figure(figsize=(13.5, 7.5))
    ax1 = plt.subplot(1, 2, 1)
    tiers = ["L1 sink-only", "L2 artifact", "L3 full chain"]
    counts = [150, 14, 30]
    colors = ["#3a6ea5", "#4e9b63", "#c97537"]
    ax1.barh(tiers, counts, color=colors)
    ax1.set_title("GT tiers used in the current benchmark assets", fontsize=16)
    ax1.set_xlabel("Binary / sample count")
    for idx, value in enumerate(counts):
        ax1.text(value + 2, idx, str(value), va="center", fontsize=12)

    ax2 = plt.subplot(1, 2, 2)
    assets = ["GT-backed", "stripped peers", "mesobench\nstripped", "no-GT", "autogen L1", "catalog total"]
    values = [44, 44, 30, 94, 108, 568]
    colors2 = ["#3a6ea5", "#5f8fc7", "#4e9b63", "#b8bec5", "#c97537", "#4d5c6c"]
    ax2.bar(range(len(assets)), values, color=colors2)
    ax2.set_title("Current canonical asset inventory", fontsize=16)
    ax2.set_xticks(range(len(assets)))
    ax2.set_xticklabels(assets, rotation=25, ha="right")
    for idx, value in enumerate(values):
        ax2.text(idx, value + 6, str(value), ha="center", fontsize=11)

    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def make_structural_results_figure(path: Path) -> None:
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13.5, 7.5))
    labels = ["Positive chains", "Channel-required", "Negative expectations"]
    done = [386, 253, 27]
    total = [386, 253, 27]
    ax1.bar(labels, total, color="#dfe6ee", label="total")
    ax1.bar(labels, done, color="#4e9b63", label="matched / satisfied")
    ax1.set_title("GT-backed structural results (44 samples)")
    ax1.legend()
    for idx, value in enumerate(done):
        ax1.text(idx, value + 5, f"{value}/{total[idx]}", ha="center", fontsize=12)
    ax1.set_ylim(0, 430)

    metrics = ["Hit rate", "Channel OK", "Negative OK", "Spurious non-drop"]
    values = [100, 100, 100, 0]
    colors = ["#4e9b63", "#4e9b63", "#4e9b63", "#bc4a48"]
    ax2.bar(metrics, values, color=colors)
    ax2.set_title("Headline interpretation")
    ax2.set_ylabel("Percent or count")
    ax2.set_ylim(0, 110)
    for idx, value in enumerate(values):
        ax2.text(idx, value + 3, "0" if idx == 3 else f"{value}%", ha="center", fontsize=12)
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def make_supervision_figure(path: Path) -> None:
    with (DOCS_DIR / "current_capability_boundary_20260310.json").open() as fh:
        summary = json.load(fh)["supervision_summary"]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13.5, 7.5))
    labels = ["queued", "reviewed", "accepted", "audit-only", "rejected"]
    values = [
        summary["queue_total"],
        summary["reviewed_total"],
        summary["accepted_total"],
        summary["audit_only_total"],
        summary["rejected_total"],
    ]
    colors = ["#5f8fc7", "#3a6ea5", "#4e9b63", "#c97537", "#bc4a48"]
    ax1.bar(labels, values, color=colors)
    ax1.set_title("Phase A.5 supervision throughput")
    for idx, value in enumerate(values):
        ax1.text(idx, value + 3, str(value), ha="center", fontsize=12)
    ax1.set_ylim(0, 170)

    kinds = list(summary["accepted_by_kind"].keys())
    kind_values = list(summary["accepted_by_kind"].values())
    kind_colors = ["#c97537", "#3a6ea5", "#bc4a48", "#4e9b63"]
    ax2.bar(kinds, kind_values, color=kind_colors)
    ax2.set_title("Accepted enrichments by artifact kind")
    for idx, value in enumerate(kind_values):
        ax2.text(idx, value + 1, str(value), ha="center", fontsize=12)
    ax2.set_ylim(0, 50)
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def make_real_firmware_figure(path: Path) -> None:
    samples = [
        {
            "name": "Contiki hello-world",
            "detected": 173,
            "chains": 65,
            "channel": 21,
            "review": 64,
        },
        {
            "name": "Zephyr CVE-2020-10065",
            "detected": 104,
            "chains": 50,
            "channel": 10,
            "review": 44,
        },
    ]
    plt.figure(figsize=(12, 7))
    x = range(len(samples))
    width = 0.18
    plt.bar([i - 1.5 * width for i in x], [s["detected"] for s in samples], width=width, label="detected labels", color="#5f8fc7")
    plt.bar([i - 0.5 * width for i in x], [s["chains"] for s in samples], width=width, label="chains", color="#3a6ea5")
    plt.bar([i + 0.5 * width for i in x], [s["channel"] for s in samples], width=width, label="channel chains", color="#4e9b63")
    plt.bar([i + 1.5 * width for i in x], [s["review"] for s in samples], width=width, label="review queue", color="#c97537")
    plt.xticks(list(x), [s["name"] for s in samples])
    plt.title("Representative no-GT real-firmware runs at stage 10")
    plt.ylabel("count")
    plt.legend()
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def make_case_studies_figure(path: Path) -> None:
    img = Image.new("RGB", (1800, 1000), PALETTE["paper"])
    draw = ImageDraw.Draw(img)
    draw.text((70, 40), "Case Studies from the Latest Full Live Run", font=FONT_H1, fill=PALETTE["navy"])
    draw.text((70, 100), "Run directory: /tmp/eval_gt_backed_full_live_20260312", font=FONT_BODY, fill=PALETTE["muted"])

    cards = [
        {
            "title": "Microbench: FreeRTOS DNS",
            "sample": "cve_2018_16525_freertos_dns",
            "fill": (221, 235, 247),
            "stats": [
                "3 chains total",
                "2 final CONFIRMED",
                "2 chains at P0",
            ],
            "takeaway": "Curated CVE chains can already survive the full pipeline and stay high-priority after review.",
        },
        {
            "title": "Parser semantics: USB host",
            "sample": "cve_2021_34259_usb_host",
            "fill": (246, 233, 236),
            "stats": [
                "3 chains total",
                "all final SUSPICIOUS",
                "all MEDIUM / P1",
            ],
            "takeaway": "Structure is present, but parser/store semantics remain conservative because destination-extent proof is still weak.",
        },
        {
            "title": "Real-style BOF: uSBS",
            "sample": "usbs_tcp_echo_client_vuln_bof",
            "fill": (255, 239, 223),
            "stats": [
                "24 chains total",
                "1 CONFIRMED / HIGH / P0",
                "19 suspicious, 4 drop",
            ],
            "takeaway": "The system can already surface one curated high-risk anchor chain out of a much noisier real-style sample.",
        },
        {
            "title": "Large real firmware: Contiki",
            "sample": "contiki_cve_2020_12140_hello_world",
            "fill": (220, 239, 231),
            "stats": [
                "65 chains total",
                "21 channel-bearing",
                "64 suspicious, 24 reviewed",
            ],
            "takeaway": "Scale and cross-context structure are already visible on large firmware, but semantic review is still conservative.",
        },
    ]

    positions = [
        (70, 180, 855, 455),
        (945, 180, 1730, 455),
        (70, 515, 855, 790),
        (945, 515, 1730, 790),
    ]
    accent_colors = [PALETTE["blue"], PALETTE["red"], PALETTE["orange"], PALETTE["green"]]

    for card, xy, accent in zip(cards, positions, accent_colors):
        _box(draw, xy, fill=card["fill"], outline=accent, radius=22, width=4)
        draw.text((xy[0] + 24, xy[1] + 20), card["title"], font=FONT_H2, fill=accent)
        draw.text((xy[0] + 24, xy[1] + 62), card["sample"], font=FONT_SMALL, fill=PALETTE["muted"])
        y = xy[1] + 105
        for stat in card["stats"]:
            draw.ellipse((xy[0] + 28, y + 8, xy[0] + 44, y + 24), fill=accent)
            y = _draw_wrapped_text(draw, (xy[0] + 56, y), stat, font=FONT_BODY, fill=PALETTE["ink"], width=xy[2] - xy[0] - 90) + 10
        draw.text((xy[0] + 24, xy[3] - 108), "Why it matters", font=FONT_SMALL, fill=accent)
        _draw_wrapped_text(
            draw,
            (xy[0] + 24, xy[3] - 80),
            card["takeaway"],
            font=FONT_BODY,
            fill=PALETTE["ink"],
            width=xy[2] - xy[0] - 48,
            line_gap=6,
        )

    draw.text(
        (70, 865),
        "These cases help separate three stories: a chain that already calibrates well, a chain that remains semantically conservative, a noisy real-style BOF sample with one strong anchor, and a large real-firmware sample where structure scales faster than semantic certainty.",
        font=FONT_BODY,
        fill=PALETTE["muted"],
    )
    img.save(path)


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.7))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*PALETTE["navy"])
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.9), Inches(12.0), Inches(0.45))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(12)
        r.font.color.rgb = RGBColor(*PALETTE["muted"])


def _add_bullets(slide, box, lines: Sequence[str], *, font_size: int = 20, color: tuple[int, int, int] = PALETTE["ink"]) -> None:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        p.bullet = True
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(*color)


def _add_text_block(slide, left: float, top: float, width: float, height: float, text: str, *, font_size: int = 18, color: tuple[int, int, int] = PALETTE["ink"], font_name: str = "Aptos") -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(*color)
    run.font.name = font_name


def _add_image(slide, path: Path, *, left: float, top: float, width: float | None = None, height: float | None = None) -> None:
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def _add_section_band(slide, text: str, *, left: float, top: float, width: float, height: float, fill: tuple[int, int, int]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*fill)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p.add_run()
    r.text = text
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*PALETTE["white"])


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "SourceAgent Progress Update", "March 12, 2026 | architecture, supervision, benchmark assets, and evidence")
    _add_section_band(slide, "Main claim", left=0.7, top=1.35, width=2.0, height=0.38, fill=PALETTE["navy"])
    _add_text_block(
        slide,
        0.75,
        1.85,
        5.2,
        2.2,
        "SourceAgent is not only a source/sink detector. It is becoming a two-phase system for monolithic firmware that recovers source, object, channel, root, and check structure, then calibrates chain risk.",
        font_size=22,
    )
    _add_section_band(slide, "What is already strong", left=0.7, top=4.25, width=2.8, height=0.38, fill=PALETTE["green"])
    box = slide.shapes.add_textbox(Inches(0.75), Inches(4.75), Inches(5.2), Inches(1.9))
    _add_bullets(
        slide,
        box,
        [
            "GT-backed structural chain recovery is already very strong.",
            "Chain-level risk GT now exists for curated anchor chains.",
            "Benchmark assets and manifests are mostly in place.",
        ],
        font_size=20,
    )
    _add_image(slide, FIG_DIR / "architecture_overview.png", left=6.2, top=1.2, width=6.6)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Target Problem and Why Structure Matters")
    left = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(5.1), Inches(5.8))
    _add_bullets(
        slide,
        left,
        [
            "Target binaries are Type-II/III monolithic firmware, not Linux-like multi-process programs.",
            "Inputs enter through MMIO, ISR-filled buffers, DMA-backed buffers, and shared SRAM objects.",
            "A sink hit alone is not enough. We need to explain how attacker-controlled data becomes the dangerous sink parameter.",
            "The core recovery targets are source, sink, object, channel, sink root, and derive/check evidence.",
        ],
        font_size=22,
    )
    _add_text_block(
        slide,
        6.2,
        1.4,
        6.3,
        4.7,
        "Phase A chain classes:\n\n- label verifier verdict: VERIFIED / PARTIAL / REJECTED / UNKNOWN\n- chain verdict: CONFIRMED / SUSPICIOUS / SAFE_OR_LOW_RISK / DROP\n- risk side-band: LOW / MEDIUM / HIGH and P0 / P1 / P2",
        font_size=21,
    )
    _add_text_block(
        slide,
        6.2,
        5.9,
        6.3,
        0.8,
        "Deterministic means reproducible fact extraction. Fail-closed means missing evidence never upgrades into a stronger claim.",
        font_size=19,
        color=PALETTE["red"],
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "End-to-End Architecture")
    _add_image(slide, FIG_DIR / "architecture_overview.png", left=0.45, top=1.0, width=12.35)
    _add_text_block(
        slide,
        0.8,
        6.45,
        12.0,
        0.5,
        "Stage 1-7 establish label-level facts. Stage 8-9 assemble structure. Phase A.5 enriches low-confidence artifacts. Phase B calibrates semantic verdict and risk.",
        font_size=18,
        color=PALETTE["muted"],
    )

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "What Each Early-Phase Artifact Means")
    _add_section_band(slide, "MemoryMap", left=0.7, top=1.2, width=1.7, height=0.35, fill=PALETTE["blue"])
    _add_text_block(slide, 0.75, 1.65, 5.7, 1.1, "Stage 1 output. A binary-specific address-space map: flash, SRAM, MMIO, system regions, base address, entry point, vector table, and ISR hypotheses.", font_size=20)
    _add_section_band(slide, "MemoryAccessIndex", left=0.7, top=3.0, width=2.4, height=0.35, fill=PALETTE["blue"])
    _add_text_block(slide, 0.75, 3.45, 5.7, 1.25, "Stage 2 output. A structured index of loads/stores with target address, provenance, function context, ISR flag, global symbols, and decompiled cache.", font_size=20)
    _add_section_band(slide, "Proposal and verifier", left=0.7, top=4.95, width=2.8, height=0.35, fill=PALETTE["orange"])
    _add_text_block(slide, 0.75, 5.4, 5.7, 1.3, "Stage 6 proposes a label for each evidence pack. Stage 7 checks label-specific obligations. The verifier sits before chain linking because it must clean local source/sink facts before they are combined.", font_size=20)
    _add_text_block(
        slide,
        7.0,
        1.35,
        5.6,
        5.2,
        "Examples of required obligations:\n\n- MMIO_READ: constant base + peripheral range\n- ISR_MMIO_READ: MMIO obligations + ISR context\n- COPY_SINK: callsite match + argument extraction\n- FORMAT_STRING_SINK: printf-like call + non-literal format\n\nIf required obligations fail, the label does not become authoritative.",
        font_size=21,
    )

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "How a Chain Is Assembled")
    _add_image(slide, FIG_DIR / "chain_assembly.png", left=0.45, top=1.0, width=12.2)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Phase A.5 LLM Supervision")
    _add_image(slide, FIG_DIR / "supervision_summary.png", left=6.35, top=1.2, width=6.1)
    left_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.2), Inches(5.3), Inches(5.9))
    _add_bullets(
        slide,
        left_box,
        [
            "Role: bounded enrichment of low-confidence sources, sinks, objects, and channels before final Phase A artifacts are frozen.",
            "Implementation position: a late Phase A enrichment loop inside stage 10: supervision queue → LLM decisions → deterministic merge gates → rebuilt artifacts.",
            "It is not a second vulnerability reviewer. It does not rewrite source reachability, object binding, channel existence, or root extraction from scratch.",
            "Workflow: build supervision_queue → ask the LLM for structured decisions → pass deterministic merge gates → enrich verified labels / objects / channels only if accepted.",
            "Current aggregate run: 148 queued, 122 reviewed, 70 accepted, 30 audit-only, 22 rejected.",
            "Acceptance is currently strongest for sources and sinks. Object/channel enrichment is wired, but still conservative.",
        ],
        font_size=20,
    )

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Benchmark Assets and GT Tiers")
    _add_image(slide, FIG_DIR / "gt_tiers_and_assets.png", left=0.55, top=1.0, width=12.1)
    _add_text_block(
        slide,
        0.8,
        6.35,
        12.0,
        0.6,
        "L1 = sink-only GT for scalable strict metrics. L2 = artifact GT for sources/objects/channels/roots/derive-checks. L3 = full chain GT for publishable end-to-end evaluation.",
        font_size=18,
        color=PALETTE["muted"],
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "How GT Is Created")
    _add_text_block(
        slide,
        0.65,
        1.25,
        5.9,
        5.8,
        "GT construction modes:\n\n1. Microbench\n- small, controlled CVE reproductions\n- artifact-complete GT\n\n2. Mesobench / GT-backed\n- realistic binaries\n- draft-then-freeze workflow\n- initial chains may be auto-promoted from live runs, then manually curated\n\n3. Risk GT\n- only curated anchor chains are labeled first\n- expected_final_verdict + expected_final_risk_band + expected_review_priority",
        font_size=20,
    )
    snippet = (
        "{\n"
        '  "chain_id": "C1_evt_overflow",\n'
        '  "expected_verdict": "CONFIRMED",\n'
        '  "expected_final_verdict": "CONFIRMED",\n'
        '  "expected_final_risk_band": "HIGH",\n'
        '  "expected_review_priority": "P0"\n'
        "}"
    )
    _add_text_block(
        slide,
        6.7,
        1.6,
        5.6,
        2.3,
        snippet,
        font_size=18,
        font_name="Courier New",
    )
    _add_text_block(
        slide,
        6.7,
        4.25,
        5.6,
        2.2,
        "Example: the HCI-over-SPI microbench sample explicitly records source, objects, sink roots, derive/check facts, full chains, and anchor-level risk GT.",
        font_size=20,
    )

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "What We Can Already Claim with High Confidence")
    _add_image(slide, FIG_DIR / "structural_results.png", left=0.55, top=1.0, width=12.0)
    _add_text_block(
        slide,
        0.8,
        6.4,
        12.0,
        0.5,
        "Interpretation: Phase A structural recovery is no longer the main bottleneck. The main open problem is verdict calibration and robustness on stripped/raw binaries.",
        font_size=18,
        color=PALETTE["muted"],
    )

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Case Studies from the Latest Full Live Run")
    _add_image(slide, FIG_DIR / "case_studies.png", left=0.4, top=0.95, width=12.5)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "What We Can and Cannot Yet Claim on Real Firmware")
    _add_image(slide, FIG_DIR / "real_firmware_scan.png", left=6.55, top=1.35, width=5.9)
    _add_bullets(
        slide,
        slide.shapes.add_textbox(Inches(0.65), Inches(1.25), Inches(5.3), Inches(5.6)),
        [
            "Stage-10 runs already scale to real monolithic firmware and produce rich artifacts, chain queues, and review worklists.",
            "Representative no-GT runs on Contiki and Zephyr show dozens of chains and many cross-context paths.",
            "This is evidence of throughput and structural richness, not yet final proof of correctness.",
            "Honest limitation: the benchmark-grade proof on broad real firmware still needs a frozen real-firmware track and more ground truth.",
        ],
        font_size=21,
    )
    _add_text_block(
        slide,
        0.7,
        6.15,
        11.7,
        0.55,
        "This is the missing claim you wanted to surface explicitly: real-firmware effectiveness is promising, but not yet proven to benchmark standard.",
        font_size=18,
        color=PALETTE["red"],
    )

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Main Innovations")
    _add_bullets(
        slide,
        slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.5), Inches(5.6)),
        [
            "Targeting Type-II/III monolithic firmware rather than Linux-like software.",
            "Recovering structure beyond source/sink labels: ChannelGraph, root-aware linking, and derive/check evidence.",
            "Using LLMs only in bounded, auditable roles: Phase A.5 supervision and Phase B semantic review cannot replace deterministic fact extraction.",
            "Building a benchmark-ready asset stack instead of a one-off detector demo.",
            "Introducing chain-level risk GT so evaluation can ask not only “was the chain found?” but also “was it calibrated to CONFIRMED / HIGH / P0?”",
        ],
        font_size=22,
    )

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "What Is Still Missing")
    _add_bullets(
        slide,
        slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.5), Inches(5.6)),
        [
            "Benchmark v1 still needs frozen dev/report splits.",
            "Raw .bin known-base and inferred-base tracks are not yet canonical benchmark tracks.",
            "Risk calibration needs a dedicated frozen subset and more patched / negative GT.",
            "Ablations exist conceptually, but are not yet frozen as benchmark presets.",
            "Suite summaries should report by format, size, execution model, and framework.",
            "The strongest remaining external criticism is real-firmware validity: we need more benchmark-grade proof, not just stress-run evidence.",
        ],
        font_size=22,
    )

    prs.save(PPT_PATH)


def build_all() -> None:
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    make_architecture_figure(FIG_DIR / "architecture_overview.png")
    make_chain_assembly_figure(FIG_DIR / "chain_assembly.png")
    make_gt_tiers_figure(FIG_DIR / "gt_tiers_and_assets.png")
    make_structural_results_figure(FIG_DIR / "structural_results.png")
    make_supervision_figure(FIG_DIR / "supervision_summary.png")
    make_case_studies_figure(FIG_DIR / "case_studies.png")
    make_real_firmware_figure(FIG_DIR / "real_firmware_scan.png")
    build_ppt()


if __name__ == "__main__":
    build_all()
